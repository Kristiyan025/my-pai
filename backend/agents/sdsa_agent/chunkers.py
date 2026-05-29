"""
Document Chunker utilities for SDSA Agent.
Provides chunking for various document types.
"""

from typing import Generator, Dict, Any, List, Optional
from dataclasses import dataclass
from enum import Enum
import re
import json


class ChunkType(str, Enum):
    TEXT = "text"
    IMAGE = "image"


@dataclass
class DocumentChunk:
    """Represents a chunk of a document."""
    content: Any  # Text string or image bytes
    chunk_type: ChunkType
    chunk_index: int
    metadata: Dict[str, Any]


# ============== Text Chunker ==============

class TextChunker:
    """
    Chunks plain text documents (txt, code files, config files).
    Uses semantic boundaries when possible.
    """
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    
    def chunk(self, content: str, metadata: Optional[Dict] = None) -> Generator[DocumentChunk, None, None]:
        """
        Chunk text content into smaller pieces.
        
        Args:
            content: Text content to chunk
            metadata: Optional metadata to include with chunks
            
        Yields:
            DocumentChunk objects
        """
        if not content:
            return
        
        base_metadata = metadata or {}
        
        # Split into paragraphs first
        paragraphs = re.split(r'\n\s*\n', content)
        
        current_chunk = ""
        chunk_index = 0
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            # If paragraph fits in current chunk, add it
            if len(current_chunk) + len(para) + 2 <= self.chunk_size:
                if current_chunk:
                    current_chunk += "\n\n" + para
                else:
                    current_chunk = para
            else:
                # Yield current chunk if not empty
                if current_chunk:
                    yield DocumentChunk(
                        content=current_chunk,
                        chunk_type=ChunkType.TEXT,
                        chunk_index=chunk_index,
                        metadata={**base_metadata, "chunk_type": "text"}
                    )
                    chunk_index += 1
                
                # Handle very long paragraphs
                if len(para) > self.chunk_size:
                    # Split by sentences
                    sentences = re.split(r'(?<=[.!?])\s+', para)
                    current_chunk = ""
                    
                    for sentence in sentences:
                        if len(current_chunk) + len(sentence) + 1 <= self.chunk_size:
                            current_chunk = current_chunk + " " + sentence if current_chunk else sentence
                        else:
                            if current_chunk:
                                yield DocumentChunk(
                                    content=current_chunk,
                                    chunk_type=ChunkType.TEXT,
                                    chunk_index=chunk_index,
                                    metadata={**base_metadata, "chunk_type": "text"}
                                )
                                chunk_index += 1
                            current_chunk = sentence
                else:
                    current_chunk = para
        
        # Yield last chunk
        if current_chunk:
            yield DocumentChunk(
                content=current_chunk,
                chunk_type=ChunkType.TEXT,
                chunk_index=chunk_index,
                metadata={**base_metadata, "chunk_type": "text"}
            )


# ============== Notebook Chunker ==============

class NotebookChunker:
    """
    Chunks Jupyter notebook files (.ipynb).
    Extracts cells with their outputs, separating text and image content.
    """
    
    def __init__(self, chunk_size: int = 1500):
        self.chunk_size = chunk_size
        self.text_chunker = TextChunker(chunk_size=chunk_size)
    
    def chunk(self, content: str, metadata: Optional[Dict] = None) -> Generator[DocumentChunk, None, None]:
        """
        Chunk notebook content.
        
        Args:
            content: JSON string of the notebook
            metadata: Optional metadata
            
        Yields:
            DocumentChunk objects (text and image chunks)
        """
        try:
            notebook = json.loads(content)
        except json.JSONDecodeError:
            # If not valid JSON, treat as text
            yield from self.text_chunker.chunk(content, metadata)
            return
        
        base_metadata = metadata or {}
        chunk_index = 0
        cells = notebook.get('cells', [])
        
        accumulated_text = []
        
        for cell_idx, cell in enumerate(cells):
            cell_type = cell.get('cell_type', 'unknown')
            source = ''.join(cell.get('source', []))
            outputs = cell.get('outputs', [])
            
            cell_metadata = {
                **base_metadata,
                "cell_index": cell_idx,
                "cell_type": cell_type
            }
            
            # Process cell source
            if cell_type in ('code', 'markdown'):
                cell_text = f"[{cell_type.upper()} Cell {cell_idx}]\n{source}"
                
                # Process outputs for code cells
                output_text = ""
                images = []
                
                for output in outputs:
                    output_type = output.get('output_type', '')
                    
                    if output_type == 'stream':
                        output_text += output.get('text', '')
                    elif output_type in ('execute_result', 'display_data'):
                        data = output.get('data', {})
                        
                        # Check for text
                        if 'text/plain' in data:
                            text_data = data['text/plain']
                            if isinstance(text_data, list):
                                text_data = ''.join(text_data)
                            output_text += str(text_data)
                        
                        # Check for images
                        for img_type in ['image/png', 'image/jpeg', 'image/gif']:
                            if img_type in data:
                                import base64
                                img_data = data[img_type]
                                if isinstance(img_data, list):
                                    img_data = ''.join(img_data)
                                images.append(base64.b64decode(img_data))
                
                if output_text:
                    cell_text += f"\n[OUTPUT]\n{output_text}"
                
                accumulated_text.append(cell_text)
                
                # Yield images separately
                for img_idx, img_data in enumerate(images):
                    # First, yield accumulated text if any
                    if accumulated_text:
                        combined_text = "\n\n".join(accumulated_text)
                        if len(combined_text) > self.chunk_size:
                            # Chunk the text
                            for sub_chunk in self.text_chunker.chunk(combined_text, cell_metadata):
                                sub_chunk.chunk_index = chunk_index
                                chunk_index += 1
                                yield sub_chunk
                        else:
                            yield DocumentChunk(
                                content=combined_text,
                                chunk_type=ChunkType.TEXT,
                                chunk_index=chunk_index,
                                metadata={**cell_metadata, "chunk_type": "text"}
                            )
                            chunk_index += 1
                        accumulated_text = []
                    
                    # Yield image chunk
                    yield DocumentChunk(
                        content=img_data,
                        chunk_type=ChunkType.IMAGE,
                        chunk_index=chunk_index,
                        metadata={
                            **cell_metadata,
                            "chunk_type": "image",
                            "image_index": img_idx
                        }
                    )
                    chunk_index += 1
                
                # Check if we should yield accumulated text
                total_len = sum(len(t) for t in accumulated_text)
                if total_len >= self.chunk_size:
                    combined_text = "\n\n".join(accumulated_text)
                    for sub_chunk in self.text_chunker.chunk(combined_text, cell_metadata):
                        sub_chunk.chunk_index = chunk_index
                        chunk_index += 1
                        yield sub_chunk
                    accumulated_text = []
        
        # Yield remaining text
        if accumulated_text:
            combined_text = "\n\n".join(accumulated_text)
            for sub_chunk in self.text_chunker.chunk(combined_text, base_metadata):
                sub_chunk.chunk_index = chunk_index
                chunk_index += 1
                yield sub_chunk


# ============== Markdown Chunker ==============

class MarkdownChunker:
    """
    Chunks Markdown files (.md, .mdx) with semantic awareness.
    Uses headers as natural boundaries.
    """
    
    def __init__(self, chunk_size: int = 1500):
        self.chunk_size = chunk_size
    
    def chunk(self, content: str, metadata: Optional[Dict] = None) -> Generator[DocumentChunk, None, None]:
        """
        Chunk markdown content using headers as boundaries.
        """
        base_metadata = metadata or {}
        
        # Split by headers
        sections = re.split(r'(^#{1,6}\s+.+$)', content, flags=re.MULTILINE)
        
        chunk_index = 0
        current_chunk = ""
        current_header = ""
        
        i = 0
        while i < len(sections):
            section = sections[i].strip()
            
            if not section:
                i += 1
                continue
            
            # Check if this is a header
            if re.match(r'^#{1,6}\s+', section):
                # If we have accumulated content, yield it
                if current_chunk:
                    yield DocumentChunk(
                        content=current_chunk,
                        chunk_type=ChunkType.TEXT,
                        chunk_index=chunk_index,
                        metadata={**base_metadata, "chunk_type": "text", "header": current_header}
                    )
                    chunk_index += 1
                    current_chunk = ""
                
                current_header = section
                current_chunk = section
            else:
                # Content section
                if len(current_chunk) + len(section) + 2 <= self.chunk_size:
                    current_chunk = f"{current_chunk}\n\n{section}" if current_chunk else section
                else:
                    # Yield current and start new
                    if current_chunk:
                        yield DocumentChunk(
                            content=current_chunk,
                            chunk_type=ChunkType.TEXT,
                            chunk_index=chunk_index,
                            metadata={**base_metadata, "chunk_type": "text", "header": current_header}
                        )
                        chunk_index += 1
                    
                    # Handle long sections
                    if len(section) > self.chunk_size:
                        text_chunker = TextChunker(chunk_size=self.chunk_size)
                        for sub_chunk in text_chunker.chunk(section, {**base_metadata, "header": current_header}):
                            sub_chunk.chunk_index = chunk_index
                            chunk_index += 1
                            yield sub_chunk
                        current_chunk = ""
                    else:
                        current_chunk = section
            
            i += 1
        
        # Yield remaining content
        if current_chunk:
            yield DocumentChunk(
                content=current_chunk,
                chunk_type=ChunkType.TEXT,
                chunk_index=chunk_index,
                metadata={**base_metadata, "chunk_type": "text", "header": current_header}
            )


# ============== Word Document Chunker ==============

class WordDocumentChunker:
    """
    Chunks Word documents (.doc, .docx).
    Extracts text and images with semantic chunking.
    """
    
    def __init__(self, chunk_size: int = 1500):
        self.chunk_size = chunk_size
        self.text_chunker = TextChunker(chunk_size=chunk_size)
    
    def chunk(self, content: bytes, metadata: Optional[Dict] = None) -> Generator[DocumentChunk, None, None]:
        """
        Chunk Word document content.
        
        Args:
            content: Raw document bytes
            metadata: Optional metadata
        """
        from io import BytesIO
        
        base_metadata = metadata or {}
        chunk_index = 0
        
        try:
            from docx import Document
            doc = Document(BytesIO(content))
            
            # Extract text from paragraphs
            text_content = []
            accumulated_text = ""
            
            for para_idx, para in enumerate(doc.paragraphs):
                text = para.text.strip()
                if not text:
                    continue
                
                # Check if this is a heading
                if para.style.name.startswith('Heading'):
                    # Yield accumulated text
                    if accumulated_text:
                        for sub_chunk in self.text_chunker.chunk(accumulated_text, base_metadata):
                            sub_chunk.chunk_index = chunk_index
                            chunk_index += 1
                            yield sub_chunk
                        accumulated_text = ""
                    
                    accumulated_text = f"[{para.style.name}] {text}"
                else:
                    if accumulated_text:
                        accumulated_text += f"\n\n{text}"
                    else:
                        accumulated_text = text
                    
                    # Check if we should yield
                    if len(accumulated_text) >= self.chunk_size:
                        for sub_chunk in self.text_chunker.chunk(accumulated_text, base_metadata):
                            sub_chunk.chunk_index = chunk_index
                            chunk_index += 1
                            yield sub_chunk
                        accumulated_text = ""
            
            # Yield remaining text
            if accumulated_text:
                for sub_chunk in self.text_chunker.chunk(accumulated_text, base_metadata):
                    sub_chunk.chunk_index = chunk_index
                    chunk_index += 1
                    yield sub_chunk
            
            # Extract images
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        image_data = rel.target_part.blob
                        yield DocumentChunk(
                            content=image_data,
                            chunk_type=ChunkType.IMAGE,
                            chunk_index=chunk_index,
                            metadata={**base_metadata, "chunk_type": "image"}
                        )
                        chunk_index += 1
                    except Exception:
                        pass
                        
        except Exception as e:
            print(f"Word document parsing failed: {e}")
            # Try to extract as plain text
            try:
                text = content.decode('utf-8', errors='ignore')
                for sub_chunk in self.text_chunker.chunk(text, base_metadata):
                    yield sub_chunk
            except Exception:
                pass


# ============== Presentation Chunker ==============

class PresentationChunker:
    """
    Chunks PowerPoint presentations (.pptx).
    Extracts text and images from slides.
    """
    
    def __init__(self, chunk_size: int = 1500):
        self.chunk_size = chunk_size
        self.text_chunker = TextChunker(chunk_size=chunk_size)
    
    def chunk(self, content: bytes, metadata: Optional[Dict] = None) -> Generator[DocumentChunk, None, None]:
        """
        Chunk PowerPoint content.
        """
        from io import BytesIO
        
        base_metadata = metadata or {}
        chunk_index = 0
        
        try:
            from pptx import Presentation
            prs = Presentation(BytesIO(content))
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_metadata = {**base_metadata, "slide_index": slide_idx}
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                    
                    # Check for images
                    if shape.shape_type == 13:  # Picture
                        try:
                            image = shape.image
                            image_data = image.blob
                            yield DocumentChunk(
                                content=image_data,
                                chunk_type=ChunkType.IMAGE,
                                chunk_index=chunk_index,
                                metadata={**slide_metadata, "chunk_type": "image"}
                            )
                            chunk_index += 1
                        except Exception:
                            pass
                
                # Yield slide text
                if slide_text:
                    combined_text = f"[Slide {slide_idx + 1}]\n" + "\n".join(slide_text)
                    
                    if len(combined_text) > self.chunk_size:
                        for sub_chunk in self.text_chunker.chunk(combined_text, slide_metadata):
                            sub_chunk.chunk_index = chunk_index
                            chunk_index += 1
                            yield sub_chunk
                    else:
                        yield DocumentChunk(
                            content=combined_text,
                            chunk_type=ChunkType.TEXT,
                            chunk_index=chunk_index,
                            metadata={**slide_metadata, "chunk_type": "text"}
                        )
                        chunk_index += 1
                        
        except Exception as e:
            print(f"Presentation parsing failed: {e}")


# ============== Document Chunker (Factory) ==============

class DocumentChunker:
    """
    Factory class that selects appropriate chunker based on file type.
    """
    
    def __init__(self, chunk_size: int = 1500):
        self.chunk_size = chunk_size
        self.text_chunker = TextChunker(chunk_size=chunk_size)
        self.notebook_chunker = NotebookChunker(chunk_size=chunk_size)
        self.markdown_chunker = MarkdownChunker(chunk_size=chunk_size)
        self.word_chunker = WordDocumentChunker(chunk_size=chunk_size)
        self.presentation_chunker = PresentationChunker(chunk_size=chunk_size)
    
    def chunk(self, content: Any, file_type: str, metadata: Optional[Dict] = None) -> Generator[DocumentChunk, None, None]:
        """
        Chunk content based on file type.
        
        Args:
            content: File content (bytes or string)
            file_type: File extension (e.g., 'txt', 'ipynb', 'docx')
            metadata: Optional metadata
        """
        file_type = file_type.lower().lstrip('.')
        base_metadata = {**(metadata or {}), "file_type": file_type}
        
        # Text-based files
        if file_type in ['txt', 'csv', 'json', 'yaml', 'yml', 'xml']:
            if isinstance(content, bytes):
                content = content.decode('utf-8', errors='ignore')
            yield from self.text_chunker.chunk(content, base_metadata)
        
        # Code files
        elif file_type in ['py', 'cpp', 'h', 'hpp', 'java', 'r', 'js', 'ts', 'html', 'css']:
            if isinstance(content, bytes):
                content = content.decode('utf-8', errors='ignore')
            yield from self.text_chunker.chunk(content, base_metadata)
        
        # Markdown
        elif file_type in ['md', 'mdx']:
            if isinstance(content, bytes):
                content = content.decode('utf-8', errors='ignore')
            yield from self.markdown_chunker.chunk(content, base_metadata)
        
        # Jupyter notebook
        elif file_type == 'ipynb':
            if isinstance(content, bytes):
                content = content.decode('utf-8', errors='ignore')
            yield from self.notebook_chunker.chunk(content, base_metadata)
        
        # Word documents
        elif file_type in ['doc', 'docx']:
            if isinstance(content, str):
                content = content.encode('utf-8')
            yield from self.word_chunker.chunk(content, base_metadata)
        
        # PowerPoint
        elif file_type == 'pptx':
            if isinstance(content, str):
                content = content.encode('utf-8')
            yield from self.presentation_chunker.chunk(content, base_metadata)
        
        # Images (single image chunk)
        elif file_type in ['png', 'jpeg', 'jpg', 'webp', 'bmp', 'gif']:
            if isinstance(content, str):
                import base64
                content = base64.b64decode(content)
            yield DocumentChunk(
                content=content,
                chunk_type=ChunkType.IMAGE,
                chunk_index=0,
                metadata={**base_metadata, "chunk_type": "image"}
            )
        
        # PDF (basic text extraction)
        elif file_type == 'pdf':
            if isinstance(content, str):
                content = content.encode('utf-8')
            yield from self._chunk_pdf(content, base_metadata)
        
        # Default: treat as text
        else:
            if isinstance(content, bytes):
                content = content.decode('utf-8', errors='ignore')
            yield from self.text_chunker.chunk(content, base_metadata)
    
    def _chunk_pdf(self, content: bytes, metadata: Dict) -> Generator[DocumentChunk, None, None]:
        """Extract text from PDF and chunk it."""
        from io import BytesIO
        chunk_index = 0
        
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(stream=content, filetype="pdf")
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text()
                
                page_metadata = {**metadata, "page_number": page_num + 1}
                
                if page_text.strip():
                    text = f"[Page {page_num + 1}]\n{page_text}"
                    
                    for sub_chunk in self.text_chunker.chunk(text, page_metadata):
                        sub_chunk.chunk_index = chunk_index
                        chunk_index += 1
                        yield sub_chunk
                
                # Extract images from page
                image_list = page.get_images()
                for img_idx, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        
                        yield DocumentChunk(
                            content=image_bytes,
                            chunk_type=ChunkType.IMAGE,
                            chunk_index=chunk_index,
                            metadata={**page_metadata, "chunk_type": "image", "image_index": img_idx}
                        )
                        chunk_index += 1
                    except Exception:
                        pass
            
            doc.close()
            
        except Exception as e:
            print(f"PDF parsing failed: {e}")
            # Fallback: try as text
            try:
                text = content.decode('utf-8', errors='ignore')
                yield from self.text_chunker.chunk(text, metadata)
            except Exception:
                pass
