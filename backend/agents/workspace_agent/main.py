"""
Workspace Agent - FastAPI application.

The system's reverse proxy to the per-user "file system" (workspace).
This agent coordinates between FileSystemService and DocumentsService
to provide file operations.
"""

import os
import sys
import asyncio
from typing import List, Optional
from datetime import datetime
import uuid
import io

from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Query
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, Response
from pydantic import BaseModel
import httpx

# Add parent to path for shared imports
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

from shared.config import settings

app = FastAPI(
    title="Workspace Agent",
    description="Manages file operations in the user's workspace",
    version="1.0.0"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ============== Configuration ==============

FILESYSTEM_SERVICE_URL = os.getenv("FILESYSTEM_SERVICE_URL", "http://filesystem-service:8001")
DOCUMENTS_SERVICE_URL = os.getenv("DOCUMENTS_SERVICE_URL", "http://documents-service:8002")
SDSA_AGENT_URL = os.getenv("SDSA_AGENT_URL", "http://sdsa-agent:8012")
USER_ID = 1  # Hardcoded single user


# ============== Request/Response Models ==============

class FileInfo(BaseModel):
    filepath: str
    uuid: str
    file_type: str
    created_at: datetime
    updated_at: datetime


class DirectoryListing(BaseModel):
    directory: str
    files: List[FileInfo]
    subdirectories: List[str]


class CopyRequest(BaseModel):
    src_paths: List[str]
    dest_dir: str


class MoveRequest(BaseModel):
    src_paths: List[str]
    dest_dir: str


class DeleteRequest(BaseModel):
    paths: List[str]


class RenameRequest(BaseModel):
    path: str
    new_name: str


class CreateDirectoryRequest(BaseModel):
    path: str


# ============== Default Directories ==============

DEFAULT_DIRECTORIES = ["/uploaded_documents"]


# ============== Startup Event ==============

@app.on_event("startup")
async def startup_event():
    """Initialize default directories on startup."""
    await asyncio.sleep(5)  # Wait for filesystem service to be ready
    
    for dir_path in DEFAULT_DIRECTORIES:
        try:
            async with httpx.AsyncClient(timeout=30.0) as client:
                # Try to create the directory
                response = await client.post(
                    f"{FILESYSTEM_SERVICE_URL}/subdirectories",
                    json={
                        "user_id": USER_ID,
                        "directory": "/",
                        "subdirectory": dir_path.strip("/")
                    }
                )
                if response.status_code == 200:
                    print(f"Created default directory: {dir_path}")
                elif response.status_code == 409:
                    print(f"Default directory already exists: {dir_path}")
                else:
                    print(f"Warning: Failed to create {dir_path}: {response.status_code}")
        except Exception as e:
            print(f"Warning: Could not create default directory {dir_path}: {e}")


# ============== Helper Functions ==============

def extract_text_from_ppt(content_bytes: bytes) -> str:
    """
    Extract text from legacy .ppt files (OLE2 compound document format).
    Uses olefile to read the PowerPoint Document stream and extract text.
    """
    import olefile
    import re
    
    try:
        ole = olefile.OleFileIO(io.BytesIO(content_bytes))
        
        # PowerPoint stores text in "PowerPoint Document" stream
        ppt_stream = None
        for stream in ['PowerPoint Document', 'PP97_DUALSTORAGE']:
            if ole.exists(stream):
                ppt_stream = ole.openstream(stream).read()
                break
        
        if not ppt_stream:
            ole.close()
            return "[Error: Could not find PowerPoint content stream in .ppt file]"
        
        # Extract text from the binary stream
        # Text in PPT is stored in records with specific headers
        # We look for TextBytesAtom (0x0FA8) and TextCharsAtom (0x0FA0) records
        extracted_texts = []
        
        # Try to extract Unicode text (UTF-16LE encoded strings)
        # Look for sequences that appear to be readable text
        text_bytes = ppt_stream
        
        # Method 1: Look for TextCharsAtom records (Unicode)
        i = 0
        while i < len(text_bytes) - 8:
            # Record header is: recVer (4 bits) + recInstance (12 bits) = 2 bytes, 
            #                   recType (2 bytes), recLen (4 bytes)
            rec_type = int.from_bytes(text_bytes[i+2:i+4], 'little')
            rec_len = int.from_bytes(text_bytes[i+4:i+8], 'little')
            
            if rec_type == 0x0FA0 and rec_len > 0 and rec_len < 100000:  # TextCharsAtom (Unicode)
                try:
                    text_data = text_bytes[i+8:i+8+rec_len]
                    text = text_data.decode('utf-16-le', errors='ignore')
                    text = text.strip()
                    if text and len(text) > 1:
                        # Filter out control characters and binary garbage
                        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
                        if text and not text.isspace():
                            extracted_texts.append(text)
                except:
                    pass
                i += 8 + rec_len
            elif rec_type == 0x0FA8 and rec_len > 0 and rec_len < 100000:  # TextBytesAtom (ASCII)
                try:
                    text_data = text_bytes[i+8:i+8+rec_len]
                    text = text_data.decode('latin-1', errors='ignore')
                    text = text.strip()
                    if text and len(text) > 1:
                        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
                        if text and not text.isspace():
                            extracted_texts.append(text)
                except:
                    pass
                i += 8 + rec_len
            else:
                i += 1
        
        ole.close()
        
        if extracted_texts:
            # Remove duplicates while preserving order
            seen = set()
            unique_texts = []
            for t in extracted_texts:
                if t not in seen:
                    seen.add(t)
                    unique_texts.append(t)
            return "\n\n".join(unique_texts)
        else:
            return "[Could not extract readable text from .ppt file. The file may be corrupted or use an unsupported format.]"
            
    except Exception as e:
        return f"[Error extracting text from .ppt file: {str(e)}]"


async def notify_sdsa_file_change(filepath: str, action: str, content: bytes = None, file_type: str = None):
    """
    Notify SDSA about file changes for embedding updates.
    Binary files (pdf, docx, pptx, etc.) are base64 encoded.
    """
    # File types that are binary and need base64 encoding
    BINARY_FILE_TYPES = {'pdf', 'docx', 'doc', 'pptx', 'ppt', 'xlsx', 'xls', 
                         'png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp', 'svg',
                         'zip', 'rar', '7z', 'tar', 'gz', 'mp3', 'mp4', 'wav',
                         'odt', 'ods', 'odp', 'epub'}
    
    try:
        async with httpx.AsyncClient(timeout=60.0) as client:
            if action == "delete":
                await client.post(
                    f"{SDSA_AGENT_URL}/documents/delete",
                    json={"user_id": USER_ID, "filepath": filepath}
                )
            elif action in ["create", "update"] and content and file_type:
                # Check if this is a binary file type
                is_binary = file_type.lower() in BINARY_FILE_TYPES if file_type else False
                
                if is_binary:
                    # Base64 encode binary content
                    import base64
                    encoded_content = base64.b64encode(content).decode('ascii')
                    await client.post(
                        f"{SDSA_AGENT_URL}/documents/index",
                        json={
                            "user_id": USER_ID,
                            "filepath": filepath,
                            "content": encoded_content,
                            "file_type": file_type,
                            "is_base64": True
                        }
                    )
                else:
                    # Text files can be sent directly
                    text_content = content.decode('utf-8', errors='ignore') if isinstance(content, bytes) else content
                    await client.post(
                        f"{SDSA_AGENT_URL}/documents/index",
                        json={
                            "user_id": USER_ID,
                            "filepath": filepath,
                            "content": text_content,
                            "file_type": file_type,
                            "is_base64": False
                        }
                    )
    except Exception as e:
        # Log but don't fail the main operation
        print(f"Warning: Failed to notify SDSA: {e}")


def normalize_path(path: str) -> str:
    """Normalize a file path."""
    if not path:
        return "/"
    if not path.startswith("/"):
        path = "/" + path
    # Remove trailing slash except for root
    if path != "/" and path.endswith("/"):
        path = path.rstrip("/")
    return path


def get_parent_directory(filepath: str) -> str:
    """Get parent directory of a path."""
    normalized = normalize_path(filepath)
    if normalized == "/":
        return "/"
    parts = normalized.rsplit("/", 1)
    return parts[0] if parts[0] else "/"


def get_filename(filepath: str) -> str:
    """Get filename from path."""
    return filepath.rsplit("/", 1)[-1]


def get_file_extension(filepath: str) -> str:
    """Get file extension without dot."""
    filename = get_filename(filepath)
    if "." in filename:
        return filename.rsplit(".", 1)[-1].lower()
    return ""


# ============== File Operations ==============

@app.get("/list", response_model=DirectoryListing)
async def list_directory(
    directory: str = Query("/", description="Directory path to list")
):
    """List files and subdirectories in a directory."""
    normalized_dir = normalize_path(directory)
    
    async with httpx.AsyncClient(timeout=30.0) as client:
        try:
            response = await client.get(
                f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}/list",
                params={"directory": normalized_dir}
            )
            response.raise_for_status()
            data = response.json()
            
            return DirectoryListing(
                directory=normalized_dir,
                files=[FileInfo(
                    filepath=f["filepath"],
                    uuid=f["uuid"],
                    file_type=f["file_type"],
                    created_at=f["created_at"],
                    updated_at=f["updated_at"]
                ) for f in data.get("files", [])],
                subdirectories=data.get("subdirectories", [])
            )
        except httpx.HTTPStatusError as e:
            raise HTTPException(status_code=e.response.status_code, detail=str(e))


@app.post("/upload")
async def upload_file(
    file: UploadFile = File(...),
    directory: str = Form("/", description="Target directory")
):
    """
    Upload a new file to the workspace.
    Returns error if file already exists at the path.
    """
    normalized_dir = normalize_path(directory)
    filepath = f"{normalized_dir}/{file.filename}" if normalized_dir != "/" else f"/{file.filename}"
    filepath = normalize_path(filepath)
    
    # Check if file already exists
    async with httpx.AsyncClient(timeout=30.0) as client:
        check_response = await client.get(
            f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
            params={"filepath": filepath}
        )
        if check_response.status_code == 200:
            raise HTTPException(status_code=409, detail="File already exists at this path")
    
    # Read file content
    content = await file.read()
    file_type = get_file_extension(file.filename)
    doc_uuid = str(uuid.uuid4())
    
    async with httpx.AsyncClient(timeout=60.0) as client:
        # Upload to DocumentsService
        files_data = {"file": (file.filename, io.BytesIO(content), file.content_type or "application/octet-stream")}
        doc_response = await client.post(
            f"{DOCUMENTS_SERVICE_URL}/documents",
            files=files_data,
            params={"custom_uuid": doc_uuid}
        )
        doc_response.raise_for_status()
        
        # Create file record in FileSystemService
        fs_response = await client.post(
            f"{FILESYSTEM_SERVICE_URL}/files",
            json={
                "user_id": USER_ID,
                "filepath": filepath,
                "uuid": doc_uuid,
                "file_type": file_type
            }
        )
        fs_response.raise_for_status()
        
        # Notify SDSA for indexing (true fire and forget - don't block response)
        asyncio.create_task(notify_sdsa_file_change(filepath, "create", content, file_type))
        
        return {
            "status": "uploaded",
            "filepath": filepath,
            "uuid": doc_uuid,
            "file_type": file_type,
            "size": len(content)
        }


class WriteContentRequest(BaseModel):
    path: str
    content: str
    encoding: str = "utf-8"


@app.post("/write")
async def write_content(request: WriteContentRequest):
    """
    Write text content to a file.
    Creates the file if it doesn't exist, overwrites if it does.
    """
    import base64 as b64_module
    
    filepath = normalize_path(request.path)
    content_bytes = request.content.encode(request.encoding)
    file_type = get_file_extension(filepath)
    
    async with httpx.AsyncClient(timeout=60.0) as client:
        # Check if file exists
        check_response = await client.get(
            f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
            params={"filepath": filepath}
        )
        
        doc_uuid = str(uuid.uuid4())
        filename = filepath.split("/")[-1]
        
        # Upload to DocumentsService
        files_data = {"file": (filename, io.BytesIO(content_bytes), "text/plain")}
        doc_response = await client.post(
            f"{DOCUMENTS_SERVICE_URL}/documents",
            files=files_data,
            params={"custom_uuid": doc_uuid}
        )
        doc_response.raise_for_status()
        
        if check_response.status_code == 200:
            # File exists - update it
            old_record = check_response.json()
            old_uuid = old_record["uuid"]
            
            # Update file record with new UUID
            await client.put(
                f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
                params={"filepath": filepath},
                json={"uuid": doc_uuid}
            )
            
            # Delete old document
            await client.delete(f"{DOCUMENTS_SERVICE_URL}/documents/{old_uuid}")
            
            # Notify SDSA
            asyncio.create_task(notify_sdsa_file_change(filepath, "update", content_bytes, file_type))
            
            return {
                "status": "updated",
                "filepath": filepath,
                "uuid": doc_uuid,
                "size": len(content_bytes)
            }
        else:
            # New file - create record
            fs_response = await client.post(
                f"{FILESYSTEM_SERVICE_URL}/files",
                json={
                    "user_id": USER_ID,
                    "filepath": filepath,
                    "uuid": doc_uuid,
                    "file_type": file_type
                }
            )
            fs_response.raise_for_status()
            
            # Notify SDSA
            asyncio.create_task(notify_sdsa_file_change(filepath, "create", content_bytes, file_type))
            
            return {
                "status": "created",
                "filepath": filepath,
                "uuid": doc_uuid,
                "size": len(content_bytes)
            }


class SaveImageRequest(BaseModel):
    path: str
    image_base64: str
    content_type: str = "image/jpeg"


@app.post("/save-image")
async def save_image(request: SaveImageRequest):
    """
    Save a base64-encoded image to the workspace.
    """
    import base64 as b64_module
    
    filepath = normalize_path(request.path)
    
    # Decode base64 image
    try:
        image_bytes = b64_module.b64decode(request.image_base64)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Invalid base64 data: {e}")
    
    file_type = get_file_extension(filepath)
    doc_uuid = str(uuid.uuid4())
    filename = filepath.split("/")[-1]
    
    async with httpx.AsyncClient(timeout=60.0) as client:
        # Check if file exists
        check_response = await client.get(
            f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
            params={"filepath": filepath}
        )
        
        # Upload to DocumentsService
        files_data = {"file": (filename, io.BytesIO(image_bytes), request.content_type)}
        doc_response = await client.post(
            f"{DOCUMENTS_SERVICE_URL}/documents",
            files=files_data,
            params={"custom_uuid": doc_uuid}
        )
        doc_response.raise_for_status()
        
        if check_response.status_code == 200:
            # File exists - update it
            old_record = check_response.json()
            old_uuid = old_record["uuid"]
            
            await client.put(
                f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
                params={"filepath": filepath},
                json={"uuid": doc_uuid}
            )
            await client.delete(f"{DOCUMENTS_SERVICE_URL}/documents/{old_uuid}")
            
            return {
                "status": "updated",
                "filepath": filepath,
                "uuid": doc_uuid,
                "size": len(image_bytes),
                "content_type": request.content_type
            }
        else:
            # New file
            fs_response = await client.post(
                f"{FILESYSTEM_SERVICE_URL}/files",
                json={
                    "user_id": USER_ID,
                    "filepath": filepath,
                    "uuid": doc_uuid,
                    "file_type": file_type
                }
            )
            fs_response.raise_for_status()
            
            return {
                "status": "created",
                "filepath": filepath,
                "uuid": doc_uuid,
                "size": len(image_bytes),
                "content_type": request.content_type
            }


@app.get("/file")
async def read_file(
    path: str = Query(..., description="File path to read")
):
    """Read/download a file from the workspace."""
    filepath = normalize_path(path)
    
    async with httpx.AsyncClient(timeout=60.0) as client:
        # Get file record from FileSystemService
        fs_response = await client.get(
            f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
            params={"filepath": filepath}
        )
        
        if fs_response.status_code == 404:
            raise HTTPException(status_code=404, detail="File not found")
        fs_response.raise_for_status()
        
        file_record = fs_response.json()
        doc_uuid = file_record["uuid"]
        
        # Get content from DocumentsService
        doc_response = await client.get(f"{DOCUMENTS_SERVICE_URL}/documents/{doc_uuid}")
        
        if doc_response.status_code == 404:
            raise HTTPException(status_code=404, detail="File content not found")
        doc_response.raise_for_status()
        
        # Return file content
        filename = get_filename(filepath)
        content_type = doc_response.headers.get("content-type", "application/octet-stream")
        
        # URL-encode non-ASCII characters for headers
        from urllib.parse import quote
        encoded_filename = quote(filename, safe='')
        encoded_filepath = quote(filepath, safe='/')
        
        return Response(
            content=doc_response.content,
            media_type=content_type,
            headers={
                "Content-Disposition": f"attachment; filename*=UTF-8''{encoded_filename}",
                "X-File-Path": encoded_filepath,
                "X-File-UUID": doc_uuid
            }
        )


@app.get("/file/text")
async def read_file_as_text(
    path: str = Query(..., description="File path to read as text")
):
    """
    Read file and extract text content for summarization/analysis.
    Supports: txt, md, json, csv, pdf, docx, pptx, ppt, xlsx, and more.
    Returns extracted text content rather than raw binary.
    """
    filepath = normalize_path(path)
    file_type = get_file_extension(filepath).lower()
    
    async with httpx.AsyncClient(timeout=60.0) as client:
        # Get file record from FileSystemService
        fs_response = await client.get(
            f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
            params={"filepath": filepath}
        )
        
        if fs_response.status_code == 404:
            raise HTTPException(status_code=404, detail="File not found")
        fs_response.raise_for_status()
        
        file_record = fs_response.json()
        doc_uuid = file_record["uuid"]
        
        # Get content from DocumentsService
        doc_response = await client.get(f"{DOCUMENTS_SERVICE_URL}/documents/{doc_uuid}")
        
        if doc_response.status_code == 404:
            raise HTTPException(status_code=404, detail="File content not found")
        doc_response.raise_for_status()
        
        content_bytes = doc_response.content
        extracted_text = ""
        
        try:
            # Handle different file types
            if file_type in ['txt', 'md', 'json', 'csv', 'xml', 'html', 'py', 'js', 'ts', 'css', 'yaml', 'yml']:
                # Plain text files - decode directly
                try:
                    extracted_text = content_bytes.decode('utf-8')
                except UnicodeDecodeError:
                    extracted_text = content_bytes.decode('latin-1')
            
            elif file_type == 'pptx':
                # PowerPoint PPTX files (Office Open XML format)
                from pptx import Presentation
                pptx_bytes = io.BytesIO(content_bytes)
                prs = Presentation(pptx_bytes)
                
                slides_text = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_texts = [f"=== Slide {slide_num} ==="]
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text.strip())
                        if shape.has_table:
                            for row in shape.table.rows:
                                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                                slide_texts.append(row_text)
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes = slide.notes_slide.notes_text_frame.text.strip()
                        if notes:
                            slide_texts.append(f"[Notes: {notes}]")
                    slides_text.append("\n".join(slide_texts))
                extracted_text = "\n\n".join(slides_text)
            
            elif file_type == 'ppt':
                # Legacy PowerPoint PPT files (OLE2 compound document format)
                extracted_text = extract_text_from_ppt(content_bytes)
            
            elif file_type == 'docx':
                # Word documents
                from docx import Document
                docx_bytes = io.BytesIO(content_bytes)
                doc = Document(docx_bytes)
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                extracted_text = "\n\n".join(paragraphs)
            
            elif file_type == 'pdf':
                # PDF files
                import fitz  # PyMuPDF
                pdf_bytes = io.BytesIO(content_bytes)
                pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
                pages_text = []
                for page_num, page in enumerate(pdf_doc, 1):
                    page_text = page.get_text()
                    if page_text.strip():
                        pages_text.append(f"=== Page {page_num} ===\n{page_text}")
                pdf_doc.close()
                extracted_text = "\n\n".join(pages_text)
            
            elif file_type in ['xlsx', 'xls']:
                # Excel files
                import pandas as pd
                excel_bytes = io.BytesIO(content_bytes)
                xls = pd.ExcelFile(excel_bytes)
                sheets_text = []
                for sheet_name in xls.sheet_names:
                    df = pd.read_excel(xls, sheet_name=sheet_name)
                    sheets_text.append(f"=== Sheet: {sheet_name} ===\n{df.to_string()}")
                extracted_text = "\n\n".join(sheets_text)
            
            elif file_type in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp']:
                # Images - return a description
                extracted_text = f"[This is an image file: {get_filename(filepath)}. Image content cannot be extracted as text. Use image analysis tools for detailed description.]"
            
            else:
                # Unknown type - try to decode as text
                try:
                    extracted_text = content_bytes.decode('utf-8')
                except UnicodeDecodeError:
                    extracted_text = f"[Binary file: {get_filename(filepath)}. Content type: {file_type}. Cannot extract text.]"
            
        except Exception as e:
            extracted_text = f"[Error extracting text from {filepath}: {str(e)}]"
        
        return {
            "filepath": filepath,
            "file_type": file_type,
            "text_content": extracted_text,
            "char_count": len(extracted_text)
        }


@app.get("/file/pptx-slides")
async def get_pptx_slides(
    path: str = Query(..., description="File path to PPTX or PPT file")
):
    """Extract slide content from a PowerPoint file (PPTX or legacy PPT)."""
    from pptx import Presentation
    from pptx.util import Inches
    import base64
    
    filepath = normalize_path(path)
    
    if not filepath.lower().endswith(('.pptx', '.ppt')):
        raise HTTPException(status_code=400, detail="File must be a PowerPoint file")
    
    is_legacy_ppt = filepath.lower().endswith('.ppt')
    
    async with httpx.AsyncClient(timeout=60.0) as client:
        # Get file record from FileSystemService
        fs_response = await client.get(
            f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
            params={"filepath": filepath}
        )
        
        if fs_response.status_code == 404:
            raise HTTPException(status_code=404, detail="File not found")
        fs_response.raise_for_status()
        
        file_record = fs_response.json()
        doc_uuid = file_record["uuid"]
        
        # Get content from DocumentsService
        doc_response = await client.get(f"{DOCUMENTS_SERVICE_URL}/documents/{doc_uuid}")
        
        if doc_response.status_code == 404:
            raise HTTPException(status_code=404, detail="File content not found")
        doc_response.raise_for_status()
        
        content_bytes = doc_response.content
        
        # Handle legacy .ppt files differently
        if is_legacy_ppt:
            try:
                extracted_text = extract_text_from_ppt(content_bytes)
                # For legacy PPT, return text as a single "slide" since we can't parse structure
                return {
                    "filename": get_filename(filepath),
                    "slide_count": 1,
                    "is_legacy_format": True,
                    "slides": [{
                        "number": 1,
                        "texts": [extracted_text] if not extracted_text.startswith("[") else [],
                        "notes": "",
                        "legacy_notice": "This is a legacy .ppt file. Slide structure cannot be parsed - showing extracted text only."
                    }]
                }
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"Failed to parse PPT: {str(e)}")
        
        # Parse PPTX (modern format)
        try:
            pptx_bytes = io.BytesIO(content_bytes)
            prs = Presentation(pptx_bytes)
            
            slides = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_data = {
                    "number": slide_num,
                    "texts": [],
                    "notes": ""
                }
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_data["texts"].append(shape.text.strip())
                    
                    # Handle tables
                    if shape.has_table:
                        table_text = []
                        for row in shape.table.rows:
                            row_cells = [cell.text.strip() for cell in row.cells]
                            table_text.append(" | ".join(row_cells))
                        if table_text:
                            slide_data["texts"].append("\n".join(table_text))
                
                # Extract notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    slide_data["notes"] = slide.notes_slide.notes_text_frame.text.strip()
                
                slides.append(slide_data)
            
            return {
                "filename": get_filename(filepath),
                "slide_count": len(prs.slides),
                "slides": slides
            }
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to parse PPTX: {str(e)}")


@app.post("/copy")
async def copy_files(request: CopyRequest):
    """
    Copy files or directories to a destination directory.
    Supports recursive copy of directories.
    """
    results = []
    errors = []
    
    async with httpx.AsyncClient(timeout=120.0) as client:
        for src_path in request.src_paths:
            src_normalized = normalize_path(src_path)
            
            try:
                # Check if source is a file
                fs_response = await client.get(
                    f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
                    params={"filepath": src_normalized}
                )
                
                if fs_response.status_code == 200:
                    # It's a file - copy it
                    file_record = fs_response.json()
                    filename = get_filename(src_normalized)
                    dest_path = normalize_path(f"{request.dest_dir}/{filename}")
                    
                    # Copy document in MinIO
                    copy_response = await client.post(
                        f"{DOCUMENTS_SERVICE_URL}/documents/{file_record['uuid']}/copy"
                    )
                    copy_response.raise_for_status()
                    new_uuid = copy_response.json()["dest_uuid"]
                    
                    # Create new file record
                    await client.post(
                        f"{FILESYSTEM_SERVICE_URL}/files",
                        json={
                            "user_id": USER_ID,
                            "filepath": dest_path,
                            "uuid": new_uuid,
                            "file_type": file_record["file_type"]
                        }
                    )
                    
                    results.append({"src": src_normalized, "dest": dest_path, "status": "copied"})
                else:
                    # It might be a directory - copy all files under it
                    all_files_response = await client.get(
                        f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}/all",
                        params={"prefix": src_normalized + "/"}
                    )
                    
                    if all_files_response.status_code == 200:
                        files = all_files_response.json()
                        for f in files:
                            rel_path = f["filepath"][len(src_normalized):]
                            dest_path = normalize_path(f"{request.dest_dir}/{get_filename(src_normalized)}{rel_path}")
                            
                            # Copy document
                            copy_response = await client.post(
                                f"{DOCUMENTS_SERVICE_URL}/documents/{f['uuid']}/copy"
                            )
                            copy_response.raise_for_status()
                            new_uuid = copy_response.json()["dest_uuid"]
                            
                            # Create file record
                            await client.post(
                                f"{FILESYSTEM_SERVICE_URL}/files",
                                json={
                                    "user_id": USER_ID,
                                    "filepath": dest_path,
                                    "uuid": new_uuid,
                                    "file_type": f["file_type"]
                                }
                            )
                            
                            results.append({"src": f["filepath"], "dest": dest_path, "status": "copied"})
                    else:
                        errors.append({"path": src_normalized, "error": "Not found"})
                        
            except Exception as e:
                errors.append({"path": src_normalized, "error": str(e)})
    
    return {"results": results, "errors": errors}


@app.post("/move")
async def move_files(request: MoveRequest):
    """
    Move files or directories to a destination directory.
    This only changes file paths in FileSystemService (no MinIO changes).
    """
    results = []
    errors = []
    
    async with httpx.AsyncClient(timeout=120.0) as client:
        for src_path in request.src_paths:
            src_normalized = normalize_path(src_path)
            
            try:
                # Check if source is a file
                fs_response = await client.get(
                    f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
                    params={"filepath": src_normalized}
                )
                
                if fs_response.status_code == 200:
                    # It's a file - move it
                    filename = get_filename(src_normalized)
                    dest_path = normalize_path(f"{request.dest_dir}/{filename}")
                    
                    # Update file path
                    await client.put(
                        f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}/rename",
                        params={"filepath": src_normalized},
                        json={"new_filepath": dest_path}
                    )
                    
                    results.append({"src": src_normalized, "dest": dest_path, "status": "moved"})
                else:
                    # It's a directory - move all files under it
                    all_files_response = await client.get(
                        f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}/all",
                        params={"prefix": src_normalized + "/"}
                    )
                    
                    if all_files_response.status_code == 200:
                        files = all_files_response.json()
                        dir_name = get_filename(src_normalized)
                        
                        for f in files:
                            rel_path = f["filepath"][len(src_normalized):]
                            dest_path = normalize_path(f"{request.dest_dir}/{dir_name}{rel_path}")
                            
                            await client.put(
                                f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}/rename",
                                params={"filepath": f["filepath"]},
                                json={"new_filepath": dest_path}
                            )
                            
                            results.append({"src": f["filepath"], "dest": dest_path, "status": "moved"})
                    else:
                        errors.append({"path": src_normalized, "error": "Not found"})
                        
            except Exception as e:
                errors.append({"path": src_normalized, "error": str(e)})
    
    return {"results": results, "errors": errors}


@app.delete("/files")
async def delete_files(request: DeleteRequest):
    """
    Delete files or directories.
    For directories, this deletes all contents recursively.
    """
    results = []
    errors = []
    
    async with httpx.AsyncClient(timeout=120.0) as client:
        for path in request.paths:
            normalized = normalize_path(path)
            
            try:
                # Check if it's a file
                fs_response = await client.get(
                    f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
                    params={"filepath": normalized}
                )
                
                if fs_response.status_code == 200:
                    # It's a file
                    file_record = fs_response.json()
                    
                    # Delete from MinIO
                    await client.delete(f"{DOCUMENTS_SERVICE_URL}/documents/{file_record['uuid']}")
                    
                    # Delete file record
                    await client.delete(
                        f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
                        params={"filepath": normalized}
                    )
                    
                    # Notify SDSA (fire and forget)
                    asyncio.create_task(notify_sdsa_file_change(normalized, "delete"))
                    
                    results.append({"path": normalized, "status": "deleted"})
                else:
                    # It's a directory - delete all files under it
                    all_files_response = await client.get(
                        f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}/all",
                        params={"prefix": normalized + "/"}
                    )
                    
                    if all_files_response.status_code == 200:
                        files = all_files_response.json()
                        
                        for f in files:
                            # Delete from MinIO
                            await client.delete(f"{DOCUMENTS_SERVICE_URL}/documents/{f['uuid']}")
                            
                            # Delete file record
                            await client.delete(
                                f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
                                params={"filepath": f["filepath"]}
                            )
                            
                            # Notify SDSA (fire and forget)
                            asyncio.create_task(notify_sdsa_file_change(f["filepath"], "delete"))
                            
                            results.append({"path": f["filepath"], "status": "deleted"})
                    else:
                        errors.append({"path": normalized, "error": "Not found"})
                        
            except Exception as e:
                errors.append({"path": normalized, "error": str(e)})
    
    return {"results": results, "errors": errors}


@app.post("/files/delete")
async def delete_files_post(request: DeleteRequest):
    return await delete_files(request)


@app.post("/rename")
async def rename_file_or_directory(request: RenameRequest):
    """Rename a file or directory."""
    normalized = normalize_path(request.path)
    parent_dir = get_parent_directory(normalized)
    new_path = normalize_path(f"{parent_dir}/{request.new_name}")
    
    async with httpx.AsyncClient(timeout=60.0) as client:
        # Check if it's a file
        fs_response = await client.get(
            f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
            params={"filepath": normalized}
        )
        
        if fs_response.status_code == 200:
            # It's a file - rename it
            await client.put(
                f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}/rename",
                params={"filepath": normalized},
                json={"new_filepath": new_path}
            )
            
            # Notify SDSA about path change (fire and forget)
            asyncio.create_task(notify_sdsa_file_change(normalized, "delete"))
            file_record = fs_response.json()
            # Would need to re-index with new path, but for simplicity just delete old
            
            return {"status": "renamed", "old_path": normalized, "new_path": new_path}
        else:
            # It's a directory - rename via subdirectory service
            parent = get_parent_directory(normalized)
            old_name = get_filename(normalized)
            
            await client.put(
                f"{FILESYSTEM_SERVICE_URL}/subdirectories/{USER_ID}/rename",
                params={
                    "directory": parent,
                    "old_name": old_name,
                    "new_name": request.new_name
                }
            )
            
            return {"status": "renamed", "old_path": normalized, "new_path": new_path}


@app.post("/directory")
async def create_directory(request: CreateDirectoryRequest):
    """Create a new empty directory."""
    normalized = normalize_path(request.path)
    parent_dir = get_parent_directory(normalized)
    dir_name = get_filename(normalized)
    
    async with httpx.AsyncClient(timeout=30.0) as client:
        response = await client.post(
            f"{FILESYSTEM_SERVICE_URL}/subdirectories",
            json={
                "user_id": USER_ID,
                "directory": parent_dir,
                "subdirectory": dir_name
            }
        )
        response.raise_for_status()
        
        return {"status": "created", "path": normalized}


@app.put("/file")
async def update_file(
    file: UploadFile = File(...),
    path: str = Form(..., description="File path to update")
):
    """
    Update a file with new content.
    Creates a new document in MinIO and updates the UUID reference.
    """
    filepath = normalize_path(path)
    
    async with httpx.AsyncClient(timeout=60.0) as client:
        # Get existing file record
        fs_response = await client.get(
            f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
            params={"filepath": filepath}
        )
        
        if fs_response.status_code == 404:
            raise HTTPException(status_code=404, detail="File not found")
        fs_response.raise_for_status()
        
        file_record = fs_response.json()
        old_uuid = file_record["uuid"]
        
        # Upload new content
        content = await file.read()
        new_uuid = str(uuid.uuid4())
        
        files_data = {"file": (file.filename, io.BytesIO(content), file.content_type or "application/octet-stream")}
        doc_response = await client.post(
            f"{DOCUMENTS_SERVICE_URL}/documents",
            files=files_data,
            params={"custom_uuid": new_uuid}
        )
        doc_response.raise_for_status()
        
        # Update file record with new UUID
        await client.put(
            f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
            params={"filepath": filepath},
            json={"uuid": new_uuid}
        )
        
        # Delete old document
        await client.delete(f"{DOCUMENTS_SERVICE_URL}/documents/{old_uuid}")
        
        # Notify SDSA (fire and forget)
        file_type = get_file_extension(filepath)
        asyncio.create_task(notify_sdsa_file_change(filepath, "update", content, file_type))
        
        return {
            "status": "updated",
            "filepath": filepath,
            "old_uuid": old_uuid,
            "new_uuid": new_uuid,
            "size": len(content)
        }


@app.get("/exists")
async def file_exists(
    path: str = Query(..., description="File path to check")
):
    """Check if a file exists."""
    filepath = normalize_path(path)
    
    async with httpx.AsyncClient(timeout=30.0) as client:
        fs_response = await client.get(
            f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
            params={"filepath": filepath}
        )
        
        return {"exists": fs_response.status_code == 200, "path": filepath}


# ============== Re-indexing ==============

@app.post("/reindex")
async def reindex_file(
    path: str = Query(..., description="File path to re-index")
):
    """Re-index a file for RAG search."""
    filepath = normalize_path(path)
    
    # Get file info from filesystem service
    async with httpx.AsyncClient(timeout=60.0) as client:
        fs_response = await client.get(
            f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}",
            params={"filepath": filepath}
        )
        
        if fs_response.status_code == 404:
            raise HTTPException(status_code=404, detail="File not found")
        
        fs_response.raise_for_status()
        file_record = fs_response.json()
        doc_uuid = file_record.get("uuid")
        file_type = file_record.get("file_type", get_file_extension(filepath))
        
        # Download content from documents service
        doc_response = await client.get(f"{DOCUMENTS_SERVICE_URL}/documents/{doc_uuid}")
        if doc_response.status_code == 404:
            raise HTTPException(status_code=404, detail="Document not found in storage")
        
        doc_response.raise_for_status()
        content = doc_response.content
        
        # Send to SDSA for re-indexing (fire and forget)
        asyncio.create_task(notify_sdsa_file_change(filepath, "create", content, file_type))
        
        return {
            "status": "reindex_triggered",
            "filepath": filepath,
            "file_type": file_type
        }


@app.post("/reindex-all")
async def reindex_all_files():
    """Re-index all files in the workspace for RAG search."""
    async with httpx.AsyncClient(timeout=120.0) as client:
        # Get all files using the correct endpoint
        fs_response = await client.get(
            f"{FILESYSTEM_SERVICE_URL}/files/{USER_ID}/all"
        )
        fs_response.raise_for_status()
        files = fs_response.json()
        
        indexed = []
        failed = []
        
        for file_record in files:
            if file_record.get("is_directory"):
                continue
                
            filepath = file_record.get("filepath")
            doc_uuid = file_record.get("uuid")
            file_type = file_record.get("file_type", "")
            
            try:
                doc_response = await client.get(f"{DOCUMENTS_SERVICE_URL}/documents/{doc_uuid}")
                if doc_response.status_code == 200:
                    content = doc_response.content
                    asyncio.create_task(notify_sdsa_file_change(filepath, "create", content, file_type))
                    indexed.append(filepath)
                else:
                    failed.append({"path": filepath, "reason": "Document not found"})
            except Exception as e:
                failed.append({"path": filepath, "reason": str(e)})
        
        return {
            "status": "complete",
            "indexed": indexed,
            "failed": failed
        }


# ============== Health Check ==============

@app.get("/health")
def health_check():
    """Health check endpoint."""
    return {"status": "healthy", "service": "workspace-agent"}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8010)
